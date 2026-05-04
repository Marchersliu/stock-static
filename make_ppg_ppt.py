#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPG涂料（张家港）有限公司
2026年第二季度党小组会议PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿（16:9宽屏）
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色
RED = RGBColor(0xC4, 0x00, 0x00)
GOLD = RGBColor(0xD4, 0xA0, 0x17)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
PPG_BLUE = RGBColor(0x00, 0x4D, 0x99)
GRAY = RGBColor(0xF5, 0xF5, 0xF5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = GOLD; circle.line.color.rgb = GOLD
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.333), Inches(1.2))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(24); p2.font.color.rgb = GOLD; p2.alignment = PP_ALIGN.CENTER
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.6))
    p3 = tb3.text_frame.paragraphs[0]; p3.text = "PPG涂料（张家港）有限公司党支部 · 2026年4月"; p3.font.size = Pt(16); p3.font.color.rgb = WHITE; p3.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = PPG_BLUE; shape.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(0.3), Inches(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.2))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(20); p2.font.color.rgb = GOLD
    return slide

def add_content_slide(prs, title, bullets, has_subtitle=False, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = RED; top_bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = RED
    if has_subtitle and subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.5))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = subtitle; p2.font.size = Pt(16); p2.font.color.rgb = RGBColor(0x66,0x66,0x66)
    start_y = 1.6 if has_subtitle else 1.2
    content = slide.shapes.add_textbox(Inches(0.6), Inches(start_y), Inches(12), Inches(5.5))
    tf = content.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0: p = tf.paragraphs[0]
        else: p = tf.add_paragraph()
        p.text = b; p.font.size = Pt(20); p.font.color.rgb = DARK; p.space_before = Pt(12); p.space_after = Pt(6)
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
    pf = footer.text_frame.paragraphs[0]; pf.text = "PPG涂料（张家港）有限公司党支部 · 机密"; pf.font.size = Pt(10); pf.font.color.rgb = RGBColor(0x99,0x99,0x99); pf.alignment = PP_ALIGN.RIGHT
    return slide

def add_two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = RED; top_bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = RED
    lt = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.5))
    plt = lt.text_frame.paragraphs[0]; plt.text = left_title; plt.font.size = Pt(22); plt.font.bold = True; plt.font.color.rgb = PPG_BLUE
    lc = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(5.5), Inches(5.0))
    tfl = lc.text_frame; tfl.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0: p = tfl.paragraphs[0]
        else: p = tfl.add_paragraph()
        p.text = item; p.font.size = Pt(18); p.font.color.rgb = DARK; p.space_before = Pt(10)
    rt = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5))
    prt = rt.text_frame.paragraphs[0]; prt.text = right_title; prt.font.size = Pt(22); prt.font.bold = True; prt.font.color.rgb = PPG_BLUE
    rc = slide.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.5), Inches(5.0))
    tfr = rc.text_frame; tfr.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0: p = tfr.paragraphs[0]
        else: p = tfr.add_paragraph()
        p.text = item; p.font.size = Pt(18); p.font.color.rgb = DARK; p.space_before = Pt(10)
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
    pf = footer.text_frame.paragraphs[0]; pf.text = "PPG涂料（张家港）有限公司党支部 · 机密"; pf.font.size = Pt(10); pf.font.color.rgb = RGBColor(0x99,0x99,0x99); pf.alignment = PP_ALIGN.RIGHT
    return slide

def add_end_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
    p = tb.text_frame.paragraphs[0]; p.text = "谢谢大家"; p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.8))
    p2 = tb2.text_frame.paragraphs[0]; p2.text = "PPG涂料（张家港）有限公司党支部"; p2.font.size = Pt(24); p2.font.color.rgb = GOLD; p2.alignment = PP_ALIGN.CENTER
    return slide

# ===== 构建PPT =====

add_title_slide(prs, "2026年第二季度党小组会议", "学习贯彻党中央精神 · 落实苏州非公党建要求 · 推动企业高质量发展")

add_content_slide(prs, "会议议程", [
    "一、学习贯彻党中央最新精神",
    "二、苏州及张家港非公党建要求",
    "三、PPG（张家港）党建实践回顾",
    "四、党建引领业务发展",
    "五、社会责任与公益实践",
    "六、第二季度工作计划"
])

add_section_slide(prs, "一、学习贯彻党中央最新精神", "落实'两个一以贯之' · 推动党建与经营深度融合")

add_content_slide(prs, "2025-2026年党中央最新指示精神", [
    "• 2025年12月中央企业负责人会议强调：坚持和加强党的全面领导",
    "• 推动党的领导融入公司治理各环节，促进党建工作与生产经营深度融合",
    "• 落实'两个一以贯之'：党的领导与公司治理相统一",
    "• 完善'双向进入、交叉任职'领导体制",
    "• 党组织前置研究讨论重大经营管理事项",
    "• 2026年是'十五五'规划开局之年，锚定高质量发展目标"
], has_subtitle=True, subtitle="坚持党的领导 · 服务中心大局 · 推动高质量发展")

add_section_slide(prs, "二、苏州及张家港非公党建要求", "外企党委领导 · 标准化规范化建设")

add_content_slide(prs, "苏州/张家港非公企业党建工作要求", [
    "• 在张家港保税区外企党委坚强领导下开展党建工作",
    "• 推进'两个覆盖'：党的组织覆盖和党的工作覆盖",
    "• 深化'四强'党支部建设：政治功能强、支部班子强、党员队伍强、作用发挥强",
    "• 落实党建标准化、规范化建设要求",
    "• 建立'堡垒指数'考评系统，实施线上提报、实时监测、评星定级",
    "• 推动党建工作与外企经营发展深度融合",
    "• 以高质量党建引领保障企业持续健康发展"
], has_subtitle=True, subtitle="外企党委领导 · 堡垒指数考评 · 深度融合经营")

add_section_slide(prs, "三、PPG（张家港）党建实践回顾", "组织建设 · 党员队伍 · 作用发挥")

add_two_col_slide(prs, "PPG（张家港）有限公司党支部概况",
    "组织建设",
    [
        "• 支部书记由管理层党员担任",
        "• 落实'双向进入、交叉任职'",
        "• 党组织活动经费保障到位",
        "• 定期召开党小组会议",
        "• 完善党员发展培养机制"
    ],
    "党员作用",
    [
        "• 党员骨干带头攻坚技术难题",
        "• 在安全生产中发挥先锋作用",
        "• 推动绿色低碳工艺改进",
        "• 带动身边同事积极向上",
        "• 争做岗位能手和技术标兵"
    ]
)

add_section_slide(prs, "四、党建引领业务发展", "安全生产 · 绿色环保 · 技术创新")

add_content_slide(prs, "以党建引领推动企业高质量发展", [
    "• 安全生产：党员带头落实安全生产责任制，实现'零事故'目标",
    "• 绿色环保：贯彻'保护并美化世界'企业理念，推进VOCs减排和清洁生产",
    "• 技术创新：党员技术骨干参与新产品研发，提升涂料核心竞争力",
    "• 质量管控：以党员示范岗带动全员质量意识提升",
    "• 降本增效：党组织牵头开展合理化建议活动，持续优化生产工艺",
    "• 人才培养：把业务骨干培养成党员，把党员培养成业务骨干"
], has_subtitle=True, subtitle="党建促经营 · 融合出实效")

add_section_slide(prs, "五、社会责任与公益实践", "多彩社区 · 教育帮扶 · 绿色环保")

add_content_slide(prs, "PPG'多彩社区'公益项目（张家港）", [
    "• 2023年9月：PPG'多彩社区'德积幼儿园项目交付",
    "   - 组织60余名员工志愿者参与",
    "   - 绘制长江三宝（鲥鱼、河豚、刀鱼）和中华鲟等本地特色图案",
    "   - 为670余名师生重塑温馨教学环境",
    "• 2015年：南沙小学'多彩社区'项目",
    "   - 170余名志愿者粉刷近千平方米操场围墙",
    "   - 捐赠总额超过4万美元",
    "• 持续践行'保护并美化世界'企业承诺",
    "• 荣获'中国企业社会责任典范奖'等多项荣誉"
], has_subtitle=True, subtitle="热心公益 · 回馈社区 · 美化世界")

add_section_slide(prs, "六、第二季度工作计划", "目标明确 · 责任落实 · 务求实效")

add_content_slide(prs, "2026年第二季度党小组工作重点", [
    "1. 深化理论学习：组织学习贯彻党的二十大及后续会议精神",
    "2. 规范组织生活：严格落实'三会一课'制度，提高组织生活质量",
    "3. 推动融合互促：围绕生产经营重点难点，设立党员攻坚项目",
    "4. 加强队伍建设：做好发展党员工作，把优秀人才吸收到党组织中",
    "5. 强化安全环保：发挥党员在安全生产和环保合规中的监督带头作用",
    "6. 持续公益实践：策划并组织至少1次社区志愿服务活动",
    "7. 迎接上级考评：对照'堡垒指数'标准，自查自纠、补齐短板"
], has_subtitle=True, subtitle="六个到位 · 全面提升")

add_end_slide(prs)

output = "/Users/hf/.kimi_openclaw/workspace/PPG张家港_2026Q2党小组会议.pptx"
prs.save(output)
print(f"PPT已保存: {output}")
print(f"总页数: {len(prs.slides)}")
